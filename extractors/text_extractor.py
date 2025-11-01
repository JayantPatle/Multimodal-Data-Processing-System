import os
import fitz 
import docx
from pptx import Presentation
from PIL import Image
import pytesseract

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def extract_text(file_path):
    """
    Extract text from supported file types:
    - PDF (.pdf)
    - Word (.docx)
    - PowerPoint (.pptx)
    - Text/Markdown (.txt, .md)
    - Images (.png, .jpg, .jpeg)
    """

    if not os.path.exists(file_path):
        return "(error: file not found)"

    # Normalize path
    file_path = file_path.strip()

    # --- PDF ---
    if file_path.lower().endswith(".pdf"):
        text = ""
        try:
            with fitz.open(file_path) as pdf:
                for page in pdf:
                    text += page.get_text("text")
            return text.strip() or "(no text found in PDF)"
        except Exception as e:
            return f"(error reading PDF: {e})"

    # --- DOCX ---
    elif file_path.lower().endswith(".docx"):
        try:
            d = docx.Document(file_path)
            return " ".join(p.text for p in d.paragraphs).strip() or "(no text found in DOCX)"
        except Exception as e:
            return f"(error reading DOCX: {e})"

    # --- PPTX ---
    elif file_path.lower().endswith(".pptx"):
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return " ".join(text).strip() or "(no text found in PPTX)"
        except Exception as e:
            return f"(error reading PPTX: {e})"

    # --- IMAGE (OCR via Tesseract) ---
    elif file_path.lower().endswith((".png", ".jpg", ".jpeg")):
        try:
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            return text.strip() or "(no text found in image)"
        except Exception as e:
            return f"(error reading image: {e})"

    # --- Plain text / Markdown ---
    elif file_path.lower().endswith((".txt", ".md")):
        try:
            with open(file_path, encoding="utf-8", errors="ignore") as f:
                return f.read().strip() or "(no text found in file)"
        except Exception as e:
            return f"(error reading text file: {e})"

    # --- Unsupported type ---
    else:
        return "(unsupported file type)"
